"""요양시설 스튜디오 화면 기록 PPT — 캡처 + 설명.

앱을 못 띄우게 되더라도 무엇을 왜 만들었는지 남기는 것이 목적이다.
화면마다 '하는 일 / 근거 / 알아둘 것' 세 덩어리로 적는다.
근거가 없는 항목은 적지 않는다(추정은 '알아둘 것'에 표시).
"""
import glob
import os

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

F = "맑은 고딕"
INK = RGBColor(0x1A, 0x1A, 0x1A)
SUB = RGBColor(0x59, 0x59, 0x59)
PURPLE = RGBColor(0x6B, 0x2A, 0x91)
PINK = RGBColor(0xC0, 0x20, 0x6B)
LINE = RGBColor(0xC7, 0xC7, 0xC7)
BG = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF7, 0xF2, 0xF8)

# 화면별 설명 — (제목, 한 줄, 하는 일, 근거, 알아둘 것)
S = {
"01": ("홈", "하루를 열면 오늘 챙길 것이 다 보인다.",
  ["오늘 근무자·등원 현황을 맨 위에",
   "스토리 줄 = 오늘 등원한 이용자. 링 색이 상태다 —\n"
   "  그라디언트는 아직 채울 칸 있음, 빨강은 주의사항 있음, 회색은 완료",
   "숫자 넷: 등원 / 결석 / 출결 미확인 / 서식 빈 칸",
   "카드로 내려가며 공지 · 기한 지난 기록 · 수급자별 기한 ·\n"
   "  손봐야 할 기기 · 오늘 주의 · 안 채운 칸 · 프로그램 · 인계"],
  ["논문 AS-IS 문제 3 — 기록이 일과 끝에 몰린다",
   "평가지표의 주기를 앱이 세어 홈에서 먼저 알린다"],
  ["인스타그램 화면 구성을 따랐다(앱바 → 스토리 → 피드).\n"
   "현장에서 '무엇부터 볼지' 고민하지 않게 하려는 것"]),

"02": ("케어 기록", "급여제공기록지를 현장에서 바로 채운다.",
  ["출결부터 찍는다(등원/결석/미확인) — 등원한 사람만 아래에 뜬다",
   "🎤 말해서 기록: 말하면 초안이 뜨고, 확인 후 저장",
   "급여시간(시작·종료·총시간 자동)과 이동서비스 차량번호",
   "3구역 12항목: 신체활동지원 · 간호 및 처치 · 기능회복훈련",
   "구역마다 특이사항과 작성자를 따로 남긴다"],
  ["시행규칙 별지 제15호서식 급여제공기록지(주·야간보호), 개정 2019.9.27",
   "서식 유의사항 반영 — 화장실은 총 횟수(기저귀는 교환 횟수),\n"
   "  목욕은 소요시간+방법, 식사는 종류+섭취량, 프로그램은 프로그램명"],
  ["자동 확정 저장이 없다. 음성이든 화면이든 사람이 눌러야 저장된다",
   "누가 언제 무엇을 음성/화면 중 무엇으로 넣었는지 따로 쌓는다(감사 기록)"]),

"03": ("인계", "다음 교대가 알아야 할 것을 기억에 의존하지 않게.",
  ["이용자 + 종류(건강 / 정서·행동 / 가족 연락 / 기기·환경 / 기타) + 내용",
   "받은 사람이 '확인함'을 누르면 누가 언제 봤는지 남는다",
   "홈과 일지에도 함께 뜬다"],
  ["논문 AS-IS 문제 2 — 정서·인지 부담, 인계가 사람 머릿속에만 있다"],
  ["음성으로 말했는데 서식 항목이 아니면(예: \"기분이 안 좋으세요\")\n"
   "케어 기록에서 '인계로 넘기기'로 보낼 수 있다"]),

"04": ("일지", "하루 기록을 모으고, 보호자 알림장까지 만든다.",
  ["이용자별로 그날 서식·인계를 시간순으로 묶어 보여준다",
   "보호자 알림장: 적어 둔 기록에서 문장을 자동으로 만든다",
   "월간 내보내기: 한 달치 급여제공기록지를 CSV 로(엑셀에서 바로 열림)"],
  ["평가 대응·공단 제출 때 한 달치를 한 번에 뽑기 위함"],
  ["보호자에게 **보내는 기능은 일부러 안 넣었다.**\n"
   "자동 발송하면 잘못 들어간 기록이 그대로 나간다.\n"
   "문장만 만들고 사람이 읽어본 뒤 문자·카톡으로 보낸다"]),

"05": ("사정·상담", "연 1회·분기·반기로 돌아오는 것들을 한 화면에.",
  ["욕구사정 9항목 — 신체·질병·인지·의사소통·영양·가족환경·\n"
   "  주관적 욕구·자원이용·구강상태",
   "상담 — 상담일자·수급자·상담직원·상담대상자(관계)·상담내용",
   "결과평가 — 목표달성·상태변화·계획반영",
   "수급자별로 기한이 지난 것을 맨 위에 붉게"],
  ["지표 30 욕구사정 연 1회 이상(대면 원칙)",
   "지표 25 상담 분기 1회 이상 — 필수기재 5개",
   "지표 44 결과평가 반기 — 결과 반영해 계획서 30일 이내 재작성"],
  ["매뉴얼이 못박은 것을 화면에 그대로 띄웠다 —\n"
   "**단순 체크는 인정 안 된다. 판단 근거를 서술해야 한다.**",
   "욕구사정 9번(구강상태)은 추정이다. 매뉴얼이 '9개'라면서 8개만 열거해\n"
   "화면에 '※ 확인 필요'를 표시해 뒀다"]),

"06": ("급여제공계획서", "급여 시작 전에 쓰고, 동의를 받아 공단에 통보한다.",
  ["머리 10칸(등급·인정번호·유효기간·급여종류·계약일·적용기간 등)",
   "목표 + 내용 표(필요영역·세부목표·필요내용·세부 제공내용·횟수·시간)",
   "종합의견 · 작성자 · 총괄 확인자 · 동의자(성명·관계·동의일)",
   "계획서가 없는 이용자를 목록으로 알려준다"],
  ["시행규칙 별지 제11호의4서식, 개정 2021.6.30",
   "지표 22 — 연 1회 이상 수립, 확인서명, 급여개시일까지 공단 통보"],
  ["동의 정보가 비면 붉게 경고한다.\n"
   "서식상 동의 없이 공단 통보를 할 수 없기 때문"]),

"07": ("대장·점검", "빠뜨리기 쉬운 12가지를, 주기를 세어 알려준다.",
  ["신체기능 프로그램(주3회) · 인지기능 프로그램(주3회) ·\n"
   "  사회적응 프로그램(월1회)",
   "차량운행일지(매일) · 실내환기 점검(매일) · 소방설비(매월)",
   "소독·감염관리(분기) · 재난훈련(반기) · 사례관리 회의(반기) ·\n"
   "  위험도 평가(반기) · 종사자 교육(매년) · 건강검진(매년)",
   "마지막 기록일로부터 며칠 지났는지 세어 기한 초과를 붉게"],
  ["2026 재가급여(주야간보호) 평가매뉴얼",
   "지표 24·25·26(프로그램) · 28(이동서비스) · 9(환기) · 16(소독) ·\n"
   "  13(소방) · 11(재난) · 29(사례관리) · 20·21(위험도) · 6(교육) · 15(검진)"],
  ["기록하는 것 자체는 어렵지 않다. **빠뜨리는 게 문제다.**\n"
   "평가에서 '이거 언제 했냐'에 답이 안 나오면 점수가 깎인다",
   "프로그램을 셋으로 나눈 이유 — 한 대장에 묶으면 사회적응(월1회)이\n"
   "신체기능(주3회) 주기에 묻힌다"]),

"08": ("이용자", "명부와 주의사항, 그리고 오늘 남은 것.",
  ["이름(가명·코드 권장) · 생년 · 자리 · 주의사항 · 메모",
   "주의사항은 붉게 — 낙상 위험, 삼킴 곤란처럼 놓치면 사고가 되는 것",
   "이용자마다 오늘 서식에서 안 채운 칸을 보여준다"],
  ["논문 AS-IS 문제 2 — 기능 수준이 다른 이용자가 한 프로그램에 섞인다"],
  ["시연 단계에서는 **가명·코드**로 등록하도록 화면에 안내한다.\n"
   "음성 인식이 이름을 구글 서버로 보내기 때문"]),

"09": ("기기 대장", "돌봄로봇·센서와 그 매뉴얼. 우리만 있는 화면.",
  ["기기명·제조사·모델·위치·상태(사용중/점검중/고장/보관/반납)",
   "기기마다 매뉴얼·설치안내·문제해결·A/S 문의처 링크를 붙인다",
   "홈에서 '손봐야 할 기기'(고장·점검중)를 따로 알린다"],
  ["논문 AS-IS 문제 4 — 기기는 들어오기만 하고 문제가 생겨도\n"
   "  공급자에게 돌아갈 길이 없다",
   "예시 8종은 논문 TO-BE 블루프린트에서 옮겼다 —\n"
   "  근력지원·이동·이승·배설·식사·목욕·모니터링·의사소통"],
  ["**민간 ERP 조사 결과 이 칸이 어디에도 없다.**\n"
   "케어포 11개 메뉴, 이지케어 12개 메뉴 어디에도 기기·로봇이 없다",
   "파일이 아니라 링크를 받는다 — 제조사가 개정하면 저절로 최신본"]),

"10": ("근무표", "누가 언제 근무했는지 남긴다.",
  ["직원별로 주간·오전·오후·야간·휴무·연차·교육 지정",
   "오늘 근무자는 홈 맨 위에 뜬다"],
  ["지표 3 인력기준 준수 · 지표 4 인력 추가배치",
   "고시 — 근무일지에 근무일자·출퇴근시간을 적어 보관해야 한다"],
  ["출퇴근 자동 기록(NFC)은 안 만들었다. 민간 ERP 영역이다"]),

"11": ("구매요청", "현장에서 떨어지는 것을 적어 둔다.",
  ["품목·분류(위생소모품/식자재/프로그램 재료/의료간호/사무/기기/기타)",
   "수량·사유 · 상태(요청 → 구매중 → 완료/보류)",
   "처리 대기와 완료를 나눠 보여준다"],
  ["평가지표가 요구하는 것은 아니다. 현장 편의 기능"],
  ["Dolbom Studio 의 구매요청서 화면을 시설용으로 줄인 것"]),

"12": ("공지", "시설 전체에 알릴 것.",
  ["누구나 등록·삭제. 홈 맨 위 카드로 뜬다"],
  [],
  ["Dolbom Studio 의 팀 공지와 같은 구조"]),

"13": ("건의", "쓰다가 불편한 것을 남긴다.",
  ["분류 4종(개선·오류·문의·기기 이상), 상태 4종(접수·진행중·완료·보류)"],
  [],
  ["현장 의견이 있어야 다음 단계를 정할 수 있다.\n"
   "Dolbom Studio 의 '개선 요청'과 같은 자리"]),

"14": ("직원·설정", "직원·PIN·백업.",
  ["직원 등록(요양보호사·사회복지사·간호(조무)사·물리작업치료사·센터장·운전원)",
   "PIN 발급 — 숫자 4~8자리. 하나라도 걸면 로그인 화면이 생긴다",
   "케어 항목 목록 확인",
   "전체 데이터 백업(zip)"],
  ["이용자 건강정보를 다루므로 잠금이 필요하다"],
  ["PIN 은 pbkdf2-sha256 해시로만 저장 — 사람이 못 읽고 재발급만 된다",
   "**시설 안 공용PC 기준의 최소 잠금이다.**\n"
   "인터넷에 공개하려면 이걸로 부족하다(계정·2단계 인증 필요)",
   "데이터는 이 PC 로컬 JSON 뿐이다. 날리면 복구가 안 되니 주기적으로 백업"]),

"15": ("다크 모드", "밝은 화면이 눈부신 사람을 위해.",
  ["사이드바 버튼으로 전환. 계정별로 기억한다"],
  [],
  ["**순백 글씨는 눈부시다**(18.7:1). 그렇다고 낮추면 시력 저하자가 못 읽는다.\n"
   "→ 본문 #DAD7D2 로 13:1 에 맞췄다",
   "인스타그램 원색은 흰 배경에서 대비가 모자란다\n"
   "(회색 3.3 · 파랑 3.2 · 빨강 3.7 — 본문 기준 4.5).\n"
   "진한 색으로 바꾸고 원색은 그라디언트 배경에만 남겼다"]),
}

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H = prs.slide_width, prs.slide_height


def base():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(s, x, y, w, h, blocks):
    """blocks = [(글, 크기, 색, 굵게, 앞여백pt)]"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for txt, size, color, bold, space in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(space)
        p.space_after = Pt(2)
        r = p.add_run(); r.text = txt
        r.font.size, r.font.name = Pt(size), F
        r.font.color.rgb, r.font.bold = color, bold
    return tb


def clean(t):
    """PPT 는 마크다운을 모른다 — ** 표시를 지우고 줄바꿈을 공백으로 편다.
    (상자 안에서 자동 줄바꿈되므로 수동 줄바꿈은 오히려 어긋난다)"""
    t = t.replace("**", "")
    return " ".join(t.split())


def label(s, x, y, w, txt, size=12, color=INK, bold=False):
    return text(s, x, y, w, .4, [(txt, size, color, bold, 0)])


# ── 표지 ──────────────────────────────────────────────────────────────
s = base()
sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.3), W, Inches(2.9))
sh.fill.solid(); sh.fill.fore_color.rgb = SOFT
sh.line.fill.background(); sh.shadow.inherit = False
label(s, .9, 2.65, 11.5, "요양시설 스튜디오", size=40, bold=True, color=PURPLE)
label(s, .9, 3.65, 11.5, "화면 기록 — 무엇을 왜 만들었나", size=19)
label(s, .9, 4.25, 11.5,
      "주간보호센터·소형 요양원용 돌봄 업무 앱 · 1단계 (2026-08 중단 시점)",
      size=13, color=SUB)
label(s, .9, 5.65, 11.5,
      "화면마다 '하는 일 / 근거 / 알아둘 것'을 적었다. "
      "근거는 법정 서식·평가지표·서비스모델 논문에서 옮긴 것이다.",
      size=12, color=SUB)
label(s, .9, 6.1, 11.5, "국립재활원 재활보조기술연구과", size=12, color=SUB)

# ── 화면별 ────────────────────────────────────────────────────────────
for f in sorted(glob.glob("care_studio/화면기록/*.png")):
    num = os.path.basename(f)[:2]
    if num not in S:
        continue
    title, one, does, basis, notes = S[num]
    s = base()
    label(s, .55, .38, 8.0, title, size=26, bold=True, color=PURPLE)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.55), Inches(.95),
                            Inches(1.4), Emu(28575))
    ln.fill.solid(); ln.fill.fore_color.rgb = PINK
    ln.line.fill.background(); ln.shadow.inherit = False
    label(s, .55, 1.06, 7.6, one, size=13, color=INK)

    # 왼쪽 캡처
    iw, ih = Image.open(f).size
    maxw, maxh = 7.5, 5.5
    sc = min(maxw / (iw / 96), maxh / (ih / 96))
    w_in, h_in = (iw / 96) * sc, (ih / 96) * sc
    pic = s.shapes.add_picture(f, Inches(.55), Inches(1.5),
                               Inches(w_in), Inches(h_in))
    pic.line.color.rgb = LINE
    pic.line.width = Pt(.75)

    # 오른쪽 설명
    x = .55 + max(w_in, 6.9) + .25
    x = min(x, 8.4)
    bw = 13.333 - x - .55
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                             Inches(1.5), Inches(bw), Inches(5.5))
    box.fill.solid(); box.fill.fore_color.rgb = BG
    box.line.color.rgb = LINE; box.line.width = Pt(1)
    box.shadow.inherit = False
    try:
        box.adjustments[0] = 0.03
    except Exception:
        pass

    # 글자 수가 많으면 크기를 한 단계 줄인다(상자 밖으로 넘치던 것 방지)
    total = sum(len(t) for t in does + basis + notes)
    fs = 10.5 if total < 380 else (9.5 if total < 520 else 8.8)
    hs = fs + 1.2
    blocks = [("하는 일", hs, PINK, True, 0)]
    for d in does:
        blocks.append(("· " + clean(d), fs, INK, False, 3))
    if basis:
        blocks.append(("근거", hs, PINK, True, 9))
        for b in basis:
            blocks.append(("· " + clean(b), fs, SUB, False, 3))
    if notes:
        blocks.append(("알아둘 것", hs, PINK, True, 9))
        for n in notes:
            blocks.append(("· " + clean(n), fs, SUB, False, 3))
    text(s, x + .22, 1.68, bw - .44, 5.2, blocks)

prs.save("요양시설스튜디오_화면기록.pptx")
print("장수:", len(prs.slides._sldIdLst))
